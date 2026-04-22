"""
Book translation module — PPTX, PDF, EPUB, DOCX.
Uses OpenAI GPT-4o-mini for translation with vision-based context detection.
"""

import base64, io, os, re, shutil, tempfile, zipfile
from pathlib import Path
from openai import OpenAI

OPENAI_KEY = os.environ.get("OPENAI_API_KEY", "")

LANGUAGES = {
    "fr": "French", "en": "English", "es": "Spanish", "de": "German",
    "it": "Italian", "pt": "Portuguese", "zh": "Chinese (Simplified)",
    "ja": "Japanese", "ko": "Korean", "ar": "Arabic", "nl": "Dutch",
    "ru": "Russian", "pl": "Polish", "tr": "Turkish", "sv": "Swedish",
}


def _get_client():
    if not OPENAI_KEY:
        raise RuntimeError("OPENAI_API_KEY not configured")
    return OpenAI(api_key=OPENAI_KEY)


def _analyze_book_context(client, image_bytes):
    """Use vision to analyze a page image and describe the book's content."""
    try:
        b64 = base64.b64encode(image_bytes).decode()
        r = client.chat.completions.create(
            model="gpt-4o-mini",
            max_tokens=150,
            messages=[
                {"role": "user", "content": [
                    {"type": "text", "text": (
                        "Describe this book page in one sentence. "
                        "What is the subject? (e.g. 'coloring book about wild animals', "
                        "'children story about a cat', 'cookbook with Italian recipes'). "
                        "Be specific and concise."
                    )},
                    {"type": "image_url", "image_url": {
                        "url": f"data:image/png;base64,{b64}",
                        "detail": "low"
                    }},
                ]},
            ],
        )
        return r.choices[0].message.content.strip()
    except Exception:
        return ""


def _translate_batch(client, texts, lang_name, book_context=""):
    """Translate a list of texts in one API call."""
    if not texts:
        return []
    numbered = "\n".join(f"[{i}] {t}" for i, t in enumerate(texts))
    context_hint = f' The book is: {book_context}.' if book_context else ''
    r = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.3,
        messages=[
            {"role": "system", "content": (
                f"Translate each numbered line to {lang_name}.{context_hint} "
                "Use the correct meaning based on context (e.g. animal names, not figurative meanings). "
                "Keep the [N] numbering. Preserve formatting, line breaks and special characters. "
                "Return ONLY the translations with their numbers, nothing else."
            )},
            {"role": "user", "content": numbered},
        ],
    )
    raw = r.choices[0].message.content.strip()
    result = {}
    for line in raw.split("\n"):
        m = re.match(r'\[(\d+)\]\s*(.*)', line)
        if m:
            result[int(m.group(1))] = m.group(2)
    return [result.get(i, texts[i]) for i in range(len(texts))]


def _translate_single(client, text, lang_name, book_context=""):
    """Translate a single text."""
    if not text.strip():
        return text
    context_hint = f' The book is: {book_context}.' if book_context else ''
    r = client.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.3,
        messages=[
            {"role": "system", "content": (
                f"Translate to {lang_name}.{context_hint} "
                "Use the correct meaning based on context. "
                "Preserve formatting and special characters. "
                "Return ONLY the translation."
            )},
            {"role": "user", "content": text},
        ],
    )
    return r.choices[0].message.content.strip()


def _batch_translate(client, all_texts, lang_name, batch_size=20, book_context=""):
    """Translate all texts in batches."""
    translated = []
    for i in range(0, len(all_texts), batch_size):
        batch = all_texts[i:i + batch_size]
        try:
            translated.extend(_translate_batch(client, batch, lang_name, book_context))
        except Exception:
            for t in batch:
                try:
                    translated.append(_translate_single(client, t, lang_name, book_context))
                except Exception:
                    translated.append(t)
    return translated


# ── PPTX ─────────────────────────────────────────────────────────────

def translate_pptx(data: bytes, lang_code: str) -> bytes:
    """Translate PPTX preserving all formatting (fonts, colors, sizes, images)."""
    from pptx import Presentation

    lang_name = LANGUAGES.get(lang_code, lang_code)
    client = _get_client()

    prs = Presentation(io.BytesIO(data))

    # Analyze up to 3 slide images for context
    book_context = ""
    slides_list = list(prs.slides)
    sample_indices = [0]
    if len(slides_list) > 4:
        sample_indices.append(len(slides_list) // 2)
    if len(slides_list) > 1:
        sample_indices.append(min(2, len(slides_list) - 1))
    descriptions = []
    for si in sample_indices:
        try:
            slide = slides_list[si]
            for shape in slide.shapes:
                if shape.shape_type == 13:  # picture
                    img_bytes = shape.image.blob
                    descriptions.append(_analyze_book_context(client, img_bytes))
                    break
        except Exception:
            pass
    if descriptions:
        book_context = " | ".join(d for d in descriptions if d)

    for slide in prs.slides:
        paragraphs = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                full = "".join(run.text for run in para.runs).strip()
                if full:
                    paragraphs.append((shape, pi, full))

        if not paragraphs:
            continue

        all_texts = [p[2] for p in paragraphs]
        translated = _batch_translate(client, all_texts, lang_name, book_context=book_context)

        for idx, (shape, pi, _orig) in enumerate(paragraphs):
            para = shape.text_frame.paragraphs[pi]
            new_text = translated[idx] if idx < len(translated) else _orig
            if not para.runs:
                continue
            para.runs[0].text = new_text
            for run in para.runs[1:]:
                run.text = ""

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF ──────────────────────────────────────────────────────────────

def _sample_bg_color(pix, bbox, scale=2):
    """Sample background color around a text bounding box."""
    sample_points = [
        (int(bbox.x0 * scale), max(0, int((bbox.y0 - 4) * scale))),       # above
        (int(bbox.x0 * scale), min(pix.height - 1, int((bbox.y1 + 4) * scale))),  # below
        (max(0, int((bbox.x0 - 4) * scale)), int(bbox.y0 * scale)),       # left
    ]
    for x, y in sample_points:
        x = max(0, min(x, pix.width - 1))
        y = max(0, min(y, pix.height - 1))
        p = pix.pixel(x, y)
        return (p[0] / 255.0, p[1] / 255.0, p[2] / 255.0)
    return (1, 1, 1)


def _map_font(span_font):
    """Map a PDF font name to a PyMuPDF built-in font."""
    f = span_font.lower()
    if "courier" in f or "mono" in f or "consol" in f:
        return "cour"
    if "times" in f or "serif" in f or "roman" in f or "garamond" in f:
        return "tiro"
    return "helv"


def _extract_page_fonts(doc, page):
    """Extract embedded fonts from a PDF page. Returns {font_name: font_bytes}."""
    font_buffers = {}
    try:
        for xref, ext, ftype, basefont, name, enc in page.get_fonts():
            try:
                _, _, _, buf = doc.extract_font(xref)
                if buf:
                    font_buffers[name] = buf
                    if basefont and basefont != name:
                        font_buffers[basefont] = buf
            except Exception:
                pass
    except Exception:
        pass
    return font_buffers


def translate_pdf(data: bytes, lang_code: str) -> bytes:
    """Translate PDF using redact & rewrite, preserving fonts and background colors."""
    import fitz

    lang_name = LANGUAGES.get(lang_code, lang_code)
    client = _get_client()

    doc = fitz.open(stream=data, filetype="pdf")

    # Analyze up to 3 pages visually to understand book context
    book_context = ""
    sample_pages = [0]
    if len(doc) > 4:
        sample_pages.append(len(doc) // 2)
    if len(doc) > 1:
        sample_pages.append(min(2, len(doc) - 1))
    descriptions = []
    for sp in sample_pages:
        try:
            pix = doc[sp].get_pixmap(matrix=fitz.Matrix(0.5, 0.5))
            descriptions.append(_analyze_book_context(client, pix.tobytes("png")))
        except Exception:
            pass
    if descriptions:
        book_context = " | ".join(d for d in descriptions if d)

    for pi in range(len(doc)):
        page = doc[pi]
        blocks = page.get_text("dict")["blocks"]
        text_items = []

        for block in blocks:
            if block["type"] != 0:
                continue
            for line in block["lines"]:
                for span in line["spans"]:
                    txt = span["text"].strip()
                    if txt:
                        text_items.append((span, txt))

        if not text_items:
            continue

        # Extract embedded fonts before redacting (redaction may remove them)
        font_buffers = _extract_page_fonts(doc, page)

        # Render page to sample background colors before redacting
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))

        all_texts = [t[1] for t in text_items]
        translated = _batch_translate(client, all_texts, lang_name, book_context=book_context)

        # Detect background color for each span, then redact
        for span, _ in text_items:
            bbox = fitz.Rect(span["bbox"])
            bg = _sample_bg_color(pix, bbox)
            page.add_redact_annot(bbox, fill=bg)
        page.apply_redactions()

        # Write translated text using original fonts
        for idx, (span, _orig) in enumerate(text_items):
            new_text = translated[idx] if idx < len(translated) else _orig
            bbox = fitz.Rect(span["bbox"])
            font_size = span["size"]
            color_int = span["color"]
            r_c = ((color_int >> 16) & 0xFF) / 255.0
            g_c = ((color_int >> 8) & 0xFF) / 255.0
            b_c = (color_int & 0xFF) / 255.0

            span_font = span.get("font", "")
            font_buf = font_buffers.get(span_font)

            # Try to use the original embedded font
            if font_buf:
                try:
                    # Measure text width with the original font
                    font_obj = fitz.Font(fontbuffer=font_buf)
                    text_width = font_obj.text_length(new_text, fontsize=font_size)
                    if text_width > bbox.width and bbox.width > 0:
                        font_size = font_size * (bbox.width / text_width) * 0.95

                    page.insert_text(
                        fitz.Point(bbox.x0, bbox.y1 - 2),
                        new_text,
                        fontbuffer=font_buf,
                        fontsize=max(font_size, 5),
                        color=(r_c, g_c, b_c),
                    )
                    continue
                except Exception:
                    pass

            # Fallback to built-in font
            font_name = _map_font(span_font)
            test_length = fitz.get_text_length(new_text, fontname=font_name, fontsize=font_size)
            if test_length > bbox.width and bbox.width > 0:
                font_size = font_size * (bbox.width / test_length) * 0.95

            page.insert_text(
                fitz.Point(bbox.x0, bbox.y1 - 2),
                new_text,
                fontname=font_name,
                fontsize=max(font_size, 5),
                color=(r_c, g_c, b_c),
            )

    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


# ── EPUB ─────────────────────────────────────────────────────────────

def translate_epub(data: bytes, lang_code: str) -> bytes:
    """Translate EPUB by parsing XHTML content files."""
    from bs4 import BeautifulSoup

    lang_name = LANGUAGES.get(lang_code, lang_code)
    client = _get_client()

    # Analyze up to 3 images from the EPUB for context
    book_context = ""
    try:
        descriptions = []
        with zipfile.ZipFile(io.BytesIO(data), 'r') as z:
            img_exts = ('.png', '.jpg', '.jpeg', '.gif', '.webp')
            img_files = [n for n in z.namelist() if n.lower().endswith(img_exts)]
            for img_name in img_files[:3]:
                try:
                    img_bytes = z.read(img_name)
                    descriptions.append(_analyze_book_context(client, img_bytes))
                except Exception:
                    pass
        if descriptions:
            book_context = " | ".join(d for d in descriptions if d)
    except Exception:
        pass

    tmp_dir = tempfile.mkdtemp(prefix="epub_tr_")
    try:
        with zipfile.ZipFile(io.BytesIO(data), 'r') as z:
            z.extractall(tmp_dir)

        xhtml_files = []
        for root, dirs, files in os.walk(tmp_dir):
            for f in files:
                if f.endswith(('.xhtml', '.html', '.htm')):
                    fp = os.path.join(root, f)
                    if 'META-INF' not in fp:
                        xhtml_files.append(fp)

        skip_tags = {'style', 'script', 'code', 'pre'}

        for xf in xhtml_files:
            with open(xf, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')

            text_nodes = []
            for el in soup.find_all(string=True):
                if el.parent.name in skip_tags:
                    continue
                txt = el.strip()
                if txt and len(txt) > 1:
                    text_nodes.append((el, txt))

            if not text_nodes:
                continue

            all_texts = [t[1] for t in text_nodes]
            translated = _batch_translate(client, all_texts, lang_name, book_context=book_context)

            for idx, (node, _orig) in enumerate(text_nodes):
                new_text = translated[idx] if idx < len(translated) else _orig
                node.replace_with(new_text)

            with open(xf, 'w', encoding='utf-8') as f:
                f.write(str(soup))

        # Update language in content.opf
        for root, dirs, files in os.walk(tmp_dir):
            for f in files:
                if f == 'content.opf':
                    opf_path = os.path.join(root, f)
                    with open(opf_path, 'r', encoding='utf-8') as fh:
                        content = fh.read()
                    content = re.sub(
                        r'(<dc:language>)[^<]*(</dc:language>)',
                        rf'\g<1>{lang_code}\g<2>',
                        content
                    )
                    with open(opf_path, 'w', encoding='utf-8') as fh:
                        fh.write(content)

        # Repack EPUB
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            mime_path = os.path.join(tmp_dir, 'mimetype')
            if os.path.exists(mime_path):
                zout.write(mime_path, 'mimetype', compress_type=zipfile.ZIP_STORED)
            for root, dirs, files in os.walk(tmp_dir):
                for f in files:
                    fp = os.path.join(root, f)
                    arcname = os.path.relpath(fp, tmp_dir)
                    if arcname == 'mimetype':
                        continue
                    zout.write(fp, arcname)

        return buf.getvalue()
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── DOCX ─────────────────────────────────────────────────────────────

def translate_docx(data: bytes, lang_code: str) -> bytes:
    """Translate DOCX preserving all formatting (fonts, colors, sizes, images)."""
    from docx import Document

    lang_name = LANGUAGES.get(lang_code, lang_code)
    client = _get_client()

    # Analyze images inside the DOCX (it's a ZIP with word/media/ images)
    book_context = ""
    try:
        descriptions = []
        with zipfile.ZipFile(io.BytesIO(data), 'r') as z:
            img_exts = ('.png', '.jpg', '.jpeg', '.gif', '.webp')
            img_files = [n for n in z.namelist() if n.lower().endswith(img_exts)]
            for img_name in img_files[:3]:
                try:
                    img_bytes = z.read(img_name)
                    descriptions.append(_analyze_book_context(client, img_bytes))
                except Exception:
                    pass
        if descriptions:
            book_context = " | ".join(d for d in descriptions if d)
    except Exception:
        pass

    doc = Document(io.BytesIO(data))

    # Collect all paragraphs (body + tables + headers + footers)
    all_paras = []

    # Body paragraphs
    for para in doc.paragraphs:
        full = "".join(run.text for run in para.runs).strip()
        if full:
            all_paras.append((para, full))

    # Table cells
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    full = "".join(run.text for run in para.runs).strip()
                    if full:
                        all_paras.append((para, full))

    # Headers & footers
    for section in doc.sections:
        for hf in [section.header, section.footer]:
            if hf and hf.is_linked_to_previous is False:
                for para in hf.paragraphs:
                    full = "".join(run.text for run in para.runs).strip()
                    if full:
                        all_paras.append((para, full))

    if not all_paras:
        return data

    # Translate all texts
    all_texts = [p[1] for p in all_paras]
    translated = _batch_translate(client, all_texts, lang_name, book_context=book_context)

    # Apply translations preserving run formatting
    for idx, (para, _orig) in enumerate(all_paras):
        new_text = translated[idx] if idx < len(translated) else _orig
        if not para.runs:
            continue
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── Dispatcher ───────────────────────────────────────────────────────

def translate_book(data: bytes, lang_code: str, content_type: str) -> bytes:
    """Route to the correct translator based on content type."""
    if 'presentation' in content_type or content_type.endswith('.pptx'):
        return translate_pptx(data, lang_code)
    elif 'pdf' in content_type:
        return translate_pdf(data, lang_code)
    elif 'epub' in content_type:
        return translate_epub(data, lang_code)
    elif 'wordprocessingml' in content_type or content_type.endswith('.docx'):
        return translate_docx(data, lang_code)
    else:
        raise ValueError(f"Unsupported content type: {content_type}")
